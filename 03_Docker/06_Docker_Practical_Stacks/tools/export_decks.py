#!/usr/bin/env python3
"""Export generated HTML decks to PDF and raster-slide PPTX.

Playwright CLI is intentionally used for every browser render so the same tool
used for LAB web evidence also proves that every slide can render in Chromium.
"""

from __future__ import annotations

import re
import shutil
import subprocess
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[1]
DECKS = [
    "Docker_Part1_Easy",
    "Docker_Part2_Intermediate",
    "Docker_Part3_Advanced",
]


def run(command: list[str]) -> None:
    print("+", " ".join(command))
    subprocess.run(command, cwd=ROOT, check=True)


def slide_count(html_path: Path) -> int:
    return len(re.findall(r'<div class="slot">', html_path.read_text(encoding="utf-8")))


def export_pdf(stem: str) -> None:
    source = (ROOT / f"{stem}.html").resolve().as_uri()
    destination = ROOT / f"{stem}.pdf"
    run([
        "npx", "playwright", "pdf",
        "--browser", "chromium",
        "--wait-for-selector", ".slide",
        "--wait-for-timeout", "250",
        source,
        str(destination),
    ])


def capture_slides(stem: str, count: int) -> Path:
    output_dir = ROOT / ".export" / stem
    if output_dir.exists():
        shutil.rmtree(output_dir)
    output_dir.mkdir(parents=True)
    uri = (ROOT / f"{stem}.html").resolve().as_uri()
    for number in range(1, count + 1):
        target = output_dir / f"slide-{number:03d}.png"
        run([
            "npx", "playwright", "screenshot",
            "--browser", "chromium",
            "--viewport-size", "1280, 720",
            "--wait-for-selector", ".slot.active",
            "--wait-for-timeout", "80",
            f"{uri}#{number}",
            str(target),
        ])
        with Image.open(target) as image:
            if image.size != (1280, 720):
                raise RuntimeError(f"Unexpected screenshot size for {target}: {image.size}")
    return output_dir


def build_pptx(stem: str, slides_dir: Path, count: int) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]
    while presentation.slides:
        slide_id = presentation.slides._sldIdLst[0]
        presentation.part.drop_rel(slide_id.rId)
        presentation.slides._sldIdLst.remove(slide_id)
    for number in range(1, count + 1):
        slide = presentation.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(slides_dir / f"slide-{number:03d}.png"),
            0,
            0,
            width=presentation.slide_width,
            height=presentation.slide_height,
        )
    presentation.core_properties.title = stem.replace("_", " ")
    presentation.core_properties.subject = "Docker Practical Stacks"
    presentation.core_properties.comments = (
        "Raster export from the self-contained HTML source using Playwright CLI."
    )
    presentation.save(ROOT / f"{stem}.pptx")


def main() -> None:
    for stem in DECKS:
        html_path = ROOT / f"{stem}.html"
        if not html_path.exists():
            raise SystemExit(f"Missing {html_path}; run tools/build_slides.py first")
        count = slide_count(html_path)
        if count < 20:
            raise RuntimeError(f"Deck {stem} is unexpectedly short: {count} slides")
        export_pdf(stem)
        slides_dir = capture_slides(stem, count)
        build_pptx(stem, slides_dir, count)
        print(f"exported {stem}: {count} slides")


if __name__ == "__main__":
    main()

