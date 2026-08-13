#!/opt/venv/bin/python3
"""Render new_Docker_Week11_Slides.html to a 90-slide PowerPoint deck."""

from __future__ import annotations

import re
import sys
from collections import deque
from hashlib import sha256
from pathlib import Path

from PIL import Image
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches


HERE = Path(__file__).resolve().parent
HTML_FILE = HERE / "new_Docker_Week11_Slides.html"
PNG_DIR = HERE / "_slides_png"
PPTX_FILE = HERE / "new_Docker_Week11_Slides.pptx"
EXPECTED_SLIDES = 90
SUPPLEMENTARY = {
    4, 8, 14, 17, 22, 27, 30, 31, 40, 47, 50, 52, 61, 65, 67, 72,
    77, 82, 84, 85, 86, 87,
}


def normalized_text(value: str) -> str:
    """Trim lines and collapse repeated blank lines without changing wording."""
    lines = [line.strip() for line in value.replace("\r\n", "\n").split("\n")]
    result: list[str] = []
    previous_blank = False
    for line in lines:
        blank = not line
        if blank and previous_blank:
            continue
        result.append(line)
        previous_blank = blank
    return "\n".join(result).strip()


def chromium_path() -> Path:
    candidates = sorted(Path("/root/.cache/ms-playwright").glob("chromium-*/chrome-linux64/chrome"))
    if not candidates:
        raise FileNotFoundError("No installed Playwright Chromium found in /root/.cache/ms-playwright")
    return candidates[-1]


def validate_rendered_pngs(png_paths: list[Path]) -> None:
    """Reject dark deck-background viewport captures and verify distinct slides.

    Slides legitimately contain dark code blocks, so only dark pixels connected
    to an image edge count as the page/deck background.
    """
    max_dark_ratio = 0.0
    for png_path in png_paths:
        with Image.open(png_path) as image:
            if image.size != (2560, 1440):
                raise RuntimeError(
                    f"Unexpected PNG dimensions for {png_path.name}: "
                    f"{image.size}, expected (2560, 1440)"
                )
            rgb = image.convert("RGB")
            width, height = rgb.size
            pixels = rgb.load()
            seen = bytearray(width * height)
            edge_dark_pixels = 0
            pending: deque[tuple[int, int]] = deque()

            def add_if_dark(x: int, y: int) -> None:
                position = y * width + x
                r, g, b = pixels[x, y]
                if not seen[position] and r < 60 and g < 60 and b < 60:
                    seen[position] = 1
                    pending.append((x, y))

            for x in range(width):
                add_if_dark(x, 0)
                add_if_dark(x, height - 1)
            for y in range(1, height - 1):
                add_if_dark(0, y)
                add_if_dark(width - 1, y)
            while pending:
                x, y = pending.popleft()
                edge_dark_pixels += 1
                for next_x, next_y in ((x - 1, y), (x + 1, y), (x, y - 1), (x, y + 1)):
                    if 0 <= next_x < width and 0 <= next_y < height:
                        add_if_dark(next_x, next_y)
            dark_ratio = edge_dark_pixels / (width * height)
        max_dark_ratio = max(max_dark_ratio, dark_ratio)
        if dark_ratio >= 0.02:
            raise RuntimeError(
                f"PNG appears to include the dark deck viewport: {png_path.name} "
                f"has {dark_ratio:.2%} dark pixels"
            )

    first_hash = sha256((PNG_DIR / "slide_001.png").read_bytes()).hexdigest()
    fiftieth_hash = sha256((PNG_DIR / "slide_050.png").read_bytes()).hexdigest()
    if first_hash == fiftieth_hash:
        raise RuntimeError("slide_001.png and slide_050.png have identical hashes")
    print(f"Maximum dark-pixel ratio: {max_dark_ratio:.2%}", flush=True)


def render_slides() -> list[dict[str, str]]:
    PNG_DIR.mkdir(exist_ok=True)
    browser_binary = chromium_path()
    print(f"Using Chromium: {browser_binary}", flush=True)
    notes: list[dict[str, str]] = []

    with sync_playwright() as playwright:
        browser = playwright.chromium.launch(headless=True, executable_path=str(browser_binary))
        page = browser.new_page(viewport={"width": 1280, "height": 720}, device_scale_factor=2)
        try:
            page.goto(HTML_FILE.as_uri(), wait_until="networkidle")
            page.wait_for_function(
                """() => Array.from(document.images).every(img =>
                    Boolean(img.getAttribute('src')) && img.complete && img.naturalWidth > 0)""",
                timeout=30000,
            )
            page.emulate_media(media="print")
            slides = page.query_selector_all("section.slide")
            count = len(slides)
            if count != EXPECTED_SLIDES:
                raise RuntimeError(f"Expected {EXPECTED_SLIDES} section.slide elements, found {count}")

            png_paths: list[Path] = []
            for index, slide in enumerate(slides):
                data = slide.evaluate(
                    """slide => {
                        const text = selector => slide.querySelector(selector)?.innerText || '';
                        return { title: text('.s-head h2'), sub: text('.s-head .sub'), body: text('.s-body') };
                    }"""
                )
                png_path = PNG_DIR / f"slide_{index + 1:03d}.png"
                slide.scroll_into_view_if_needed()
                slide.screenshot(path=str(png_path), animations="disabled")
                png_paths.append(png_path)
                notes.append({key: normalized_text(value) for key, value in data.items()})
                print(f"Rendered {index + 1:02d}/{count}: {png_path.name}", flush=True)
            validate_rendered_pngs(png_paths)
        finally:
            browser.close()
    return notes


def make_presentation(notes: list[dict[str, str]]) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]

    for index, note in enumerate(notes, start=1):
        slide = presentation.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(PNG_DIR / f"slide_{index:03d}.png"), 0, 0,
            width=presentation.slide_width, height=presentation.slide_height,
        )
        note_lines: list[str] = []
        if index in SUPPLEMENTARY:
            note_lines.append("[สไลด์เสริม - ไม่อยู่ในชุดหลัก 68 ใบ]")
        note_lines.append(note["title"])
        if note["sub"]:
            note_lines.append(note["sub"])
        if note["body"]:
            note_lines.append(note["body"])
        slide.notes_slide.notes_text_frame.text = "\n".join(line for line in note_lines if line)
        print(f"Added PPTX slide {index:02d}/{len(notes)}", flush=True)

    presentation.save(PPTX_FILE)


def main() -> int:
    if not HTML_FILE.is_file():
        raise FileNotFoundError(f"Input not found: {HTML_FILE}")
    print(f"Rendering {HTML_FILE.name}...", flush=True)
    notes = render_slides()
    if len(notes) != EXPECTED_SLIDES:
        raise RuntimeError(f"Rendered {len(notes)} slides, expected {EXPECTED_SLIDES}")
    print("Creating PowerPoint...", flush=True)
    make_presentation(notes)
    check = Presentation(PPTX_FILE)
    if len(check.slides) != EXPECTED_SLIDES:
        raise RuntimeError(f"PPTX verification failed: {len(check.slides)} slides")
    print(f"Done: {len(check.slides)} slides; {PPTX_FILE.stat().st_size:,} bytes; {PPTX_FILE}", flush=True)
    return 0


if __name__ == "__main__":
    try:
        raise SystemExit(main())
    except Exception as error:
        print(f"ERROR: {error}", file=sys.stderr, flush=True)
        raise
